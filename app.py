import os
import re
import json
import hashlib
import secrets
from functools import lru_cache
from pathlib import Path
from typing import List, Literal

import asyncio
import torch
from fastapi import FastAPI, Request, Form
from fastapi.responses import HTMLResponse, JSONResponse, RedirectResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from starlette.middleware.sessions import SessionMiddleware
from pydantic import BaseModel, Field
from transformers import AutoModelForCausalLM, AutoTokenizer
import io
import threading
import time
from typing import Optional
from fastapi import UploadFile, File

# ── Config ────────────────────────────────────────────────────────────────────
APP_NAME       = "VedAI"
MODEL_ID       = os.getenv("MODEL_ID", "Qwen/Qwen2.5-1.5B-Instruct")
MAX_NEW_TOKENS = int(os.getenv("MAX_NEW_TOKENS", "1024"))
TEMPERATURE    = float(os.getenv("TEMPERATURE", "0.35"))
TOP_P          = float(os.getenv("TOP_P", "0.90"))
MAX_HISTORY    = int(os.getenv("MAX_HISTORY_MESSAGES", "10"))
DOC_TTL_MIN    = int(os.getenv("DOC_TTL_MINUTES", "30"))
MAX_DOC_CHARS  = int(os.getenv("MAX_DOC_CHARS", "10000"))
SHARE_TTL_HOURS = int(os.getenv("SHARE_TTL_HOURS", "24"))

# ── Admin credentials ─────────────────────────────────────────────────────────
ADMIN_USERNAME      = os.getenv("ADMIN_USERNAME", "admin")
ADMIN_PASSWORD_HASH = os.getenv("ADMIN_PASSWORD_HASH", "")
_ADMIN_PW_PLAIN     = os.getenv("ADMIN_PASSWORD", "Admin@987#hf")
SESSION_SECRET      = os.getenv("SECRET_KEY", secrets.token_hex(32))

USERS_FILE = Path(os.getenv("USERS_FILE", "/tmp/vedai_users.json"))

# ── Document store ────────────────────────────────────────────────────────────
_doc_store: dict = {}
_doc_lock = threading.Lock()

def _cleanup_docs():
    while True:
        time.sleep(60)
        now = time.time()
        with _doc_lock:
            for sid in list(_doc_store):
                for did in list(_doc_store[sid]):
                    d = _doc_store[sid][did]
                    if now - d["uploaded_at"] > d["ttl_minutes"] * 60:
                        del _doc_store[sid][did]
                if not _doc_store[sid]:
                    del _doc_store[sid]

threading.Thread(target=_cleanup_docs, daemon=True).start()

def store_doc(sid, filename, text, ttl=DOC_TTL_MIN):
    did = secrets.token_hex(8)
    with _doc_lock:
        _doc_store.setdefault(sid, {})[did] = {
            "filename": filename, "text": text,
            "uploaded_at": time.time(), "ttl_minutes": ttl, "chars": len(text),
        }
    return did

def list_docs(sid):
    now = time.time()
    out = []
    with _doc_lock:
        for did, d in list(_doc_store.get(sid, {}).items()):
            rem = d["ttl_minutes"] * 60 - (now - d["uploaded_at"])
            if rem > 0:
                out.append({"id": did, "filename": d["filename"],
                             "chars": d["chars"], "expires_in_min": round(rem/60, 1)})
    return out

def get_doc_context(sid):
    now = time.time()
    parts = []
    with _doc_lock:
        for did, d in list(_doc_store.get(sid, {}).items()):
            if now - d["uploaded_at"] <= d["ttl_minutes"] * 60:
                parts.append(f"[Document: {d['filename']}]\n{d['text'][:MAX_DOC_CHARS]}\n[End: {d['filename']}]")
    return "\n\n".join(parts)

def delete_doc(sid, did):
    with _doc_lock:
        return bool(_doc_store.get(sid, {}).pop(did, None))

# ── Document parsers ──────────────────────────────────────────────────────────
def parse_pdf(data):
    try:
        from pypdf import PdfReader
        r = PdfReader(io.BytesIO(data))
        return "\n\n".join(p.extract_text() or "" for p in r.pages).strip()
    except Exception as e: return f"[PDF error: {e}]"

def parse_docx(data):
    try:
        import docx
        doc = docx.Document(io.BytesIO(data))
        lines = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        for t in doc.tables:
            for row in t.rows:
                lines.append(" | ".join(c.text.strip() for c in row.cells if c.text.strip()))
        return "\n\n".join(lines)
    except Exception as e: return f"[DOCX error: {e}]"

def parse_pptx(data):
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        slides = []
        for i, s in enumerate(prs.slides, 1):
            texts = [sh.text.strip() for sh in s.shapes if hasattr(sh, "text") and sh.text.strip()]
            if texts: slides.append(f"[Slide {i}]\n" + "\n".join(texts))
        return "\n\n".join(slides)
    except Exception as e: return f"[PPTX error: {e}]"

def parse_txt(data):
    for enc in ("utf-8", "latin-1"):
        try: return data.decode(enc)
        except: pass
    return "[Decode error]"

def parse_image(data: bytes, filename: str = "") -> str:
    ocr_text = ""
    try:
        import pytesseract
        from PIL import Image as PILImage
        img = PILImage.open(io.BytesIO(data))
        ocr_text = pytesseract.image_to_string(img).strip()
    except Exception:
        pass
    try:
        from PIL import Image as PILImage
        img = PILImage.open(io.BytesIO(data))
        w, h = img.size
        fmt = img.format or Path(filename).suffix.upper().lstrip(".")
        meta = f"[Image: {filename} | {fmt} | {w}×{h}px]"
    except Exception:
        meta = f"[Image: {filename}]"
    if ocr_text:
        return f"{meta}\nOCR extracted text:\n{ocr_text}"
    return (f"{meta}\nThe user has shared this image. "
            f"Acknowledge it and describe or answer questions about it based on the filename and context.")

DOC_PARSERS = {
    ".pdf": parse_pdf, ".docx": parse_docx, ".doc": parse_docx,
    ".pptx": parse_pptx, ".ppt": parse_pptx, ".txt": parse_txt, ".md": parse_txt,
    ".jpg": lambda d: parse_image(d, "image.jpg"),
    ".jpeg": lambda d: parse_image(d, "image.jpeg"),
    ".png": lambda d: parse_image(d, "image.png"),
    ".gif": lambda d: parse_image(d, "image.gif"),
    ".webp": lambda d: parse_image(d, "image.webp"),
}

# ── Abbreviation dictionary (million-scale: 1000+ entries across all domains) ─
_ABBREV = {
    # ── Indian Research & Science ──────────────────────────────────────────────
    "vecc":"Variable Energy Cyclotron Centre","barc":"Bhabha Atomic Research Centre",
    "isro":"Indian Space Research Organisation","drdo":"Defence Research and Development Organisation",
    "iit":"Indian Institute of Technology","iim":"Indian Institute of Management",
    "iisc":"Indian Institute of Science","nit":"National Institute of Technology",
    "bits":"Birla Institute of Technology and Science","iiser":"Indian Institute of Science Education and Research",
    "tifr":"Tata Institute of Fundamental Research","prl":"Physical Research Laboratory",
    "iucaa":"Inter-University Centre for Astronomy and Astrophysics",
    "ncbs":"National Centre for Biological Sciences","csir":"Council of Scientific Industrial Research",
    "dae":"Department of Atomic Energy","dst":"Department of Science and Technology",
    "dbt":"Department of Biotechnology","icar":"Indian Council of Agricultural Research",
    "icmr":"Indian Council of Medical Research","nstc":"National Science and Technology Council",
    "aries":"Aryabhatta Research Institute of Observational Sciences",
    "iia":"Indian Institute of Astrophysics","rri":"Raman Research Institute",
    "iop":"Institute of Physics","imsc":"Institute of Mathematical Sciences",
    "tcs":"Tata Consultancy Services","wipro":"Western India Products",
    "infosys":"Information Systems","hcl":"Hindustan Computers Limited",
    "npcil":"Nuclear Power Corporation of India Limited",
    "igcar":"Indira Gandhi Centre for Atomic Research",
    "rrcat":"Raja Ramanna Centre for Advanced Technology",
    "sac":"Space Applications Centre","nrsc":"National Remote Sensing Centre",
    "sam":"Scientific Analysis Mission","vssc":"Vikram Sarabhai Space Centre",
    "lpsc":"Liquid Propulsion Systems Centre","ursc":"U R Rao Satellite Centre",
    "iist":"Indian Institute of Space Science and Technology",
    # ── Global Science & Space ─────────────────────────────────────────────────
    "nasa":"National Aeronautics and Space Administration",
    "esa":"European Space Agency","cern":"European Organisation for Nuclear Research",
    "nih":"National Institutes of Health","nsf":"National Science Foundation",
    "cnrs":"Centre National de la Recherche Scientifique",
    "mpg":"Max Planck Gesellschaft","dkfz":"German Cancer Research Centre",
    "jaxa":"Japan Aerospace Exploration Agency","roscosmos":"Russian State Space Corporation",
    "cnsa":"China National Space Administration",
    "noaa":"National Oceanic and Atmospheric Administration",
    "nist":"National Institute of Standards and Technology",
    "doe":"Department of Energy","darpa":"Defence Advanced Research Projects Agency",
    "arpa":"Advanced Research Projects Agency","iana":"Internet Assigned Numbers Authority",
    # ── Medicine & Health ──────────────────────────────────────────────────────
    "who":"World Health Organisation","aiims":"All India Institute of Medical Sciences",
    "dna":"deoxyribonucleic acid","rna":"ribonucleic acid","mrna":"messenger ribonucleic acid",
    "atp":"adenosine triphosphate","ecg":"electrocardiogram","eeg":"electroencephalogram",
    "mri":"magnetic resonance imaging","ct":"computed tomography","pet":"positron emission tomography",
    "bmi":"body mass index","icu":"intensive care unit","ot":"operation theatre",
    "bp":"blood pressure","hr":"heart rate","hiv":"human immunodeficiency virus",
    "aids":"acquired immunodeficiency syndrome","covid":"coronavirus disease",
    "sars":"severe acute respiratory syndrome","mers":"Middle East respiratory syndrome",
    "tb":"tuberculosis","hbp":"high blood pressure","ckd":"chronic kidney disease",
    "cad":"coronary artery disease","copd":"chronic obstructive pulmonary disease",
    "cvd":"cardiovascular disease","ncd":"non-communicable disease",
    "ors":"oral rehydration solution","otp":"oral therapeutic programme",
    "anc":"antenatal care","pnc":"postnatal care","icd":"international classification of diseases",
    "who":"World Health Organisation","picu":"paediatric intensive care unit",
    "nicu":"neonatal intensive care unit","ed":"emergency department",
    "cpr":"cardiopulmonary resuscitation","aed":"automated external defibrillator",
    "ivf":"in vitro fertilisation","iui":"intrauterine insemination",
    "hsg":"hysterosalpingography","pcod":"polycystic ovarian disease",
    "pcos":"polycystic ovarian syndrome",
    # ── Physics & Engineering ──────────────────────────────────────────────────
    "lhc":"Large Hadron Collider","cms":"Compact Muon Solenoid",
    "atlas":"A Toroidal LHC Apparatus","alice":"A Large Ion Collider Experiment",
    "ligo":"Laser Interferometer Gravitational-Wave Observatory",
    "iter":"International Thermonuclear Experimental Reactor",
    "iter":"International Thermonuclear Experimental Reactor",
    "iaea":"International Atomic Energy Agency",
    "npt":"Nuclear Non-Proliferation Treaty","ctbt":"Comprehensive Nuclear Test Ban Treaty",
    "ev":"electron volt","gev":"giga electron volt","tev":"tera electron volt",
    "pev":"peta electron volt","kev":"kilo electron volt","mev":"mega electron volt",
    "dc":"direct current","ac":"alternating current","rf":"radio frequency",
    "em":"electromagnetic","uv":"ultraviolet","ir":"infrared",
    "led":"light emitting diode","laser":"light amplification by stimulated emission of radiation",
    "maser":"microwave amplification by stimulated emission of radiation",
    "sar":"synthetic aperture radar","lidar":"light detection and ranging",
    "radar":"radio detection and ranging","sonar":"sound navigation and ranging",
    # ── Chemistry ─────────────────────────────────────────────────────────────
    "iupac":"International Union of Pure and Applied Chemistry",
    "gcms":"gas chromatography mass spectrometry","hplc":"high performance liquid chromatography",
    "nmr":"nuclear magnetic resonance","esr":"electron spin resonance",
    "ir":"infrared spectroscopy","xrd":"x-ray diffraction",
    "xrf":"x-ray fluorescence","sem":"scanning electron microscopy",
    "tem":"transmission electron microscopy","edx":"energy dispersive x-ray",
    "ato":"atomic force microscopy","stm":"scanning tunnelling microscopy",
    "tga":"thermogravimetric analysis","dsc":"differential scanning calorimetry",
    "ph":"potential of hydrogen","ppm":"parts per million","ppb":"parts per billion",
    # ── Computer Science & Tech ────────────────────────────────────────────────
    "ai":"artificial intelligence","ml":"machine learning","dl":"deep learning",
    "nlp":"natural language processing","cv":"computer vision",
    "llm":"large language model","gpt":"generative pre-trained transformer",
    "bert":"bidirectional encoder representations from transformers",
    "vit":"vision transformer","gan":"generative adversarial network",
    "vae":"variational autoencoder","cnn":"convolutional neural network",
    "rnn":"recurrent neural network","lstm":"long short-term memory",
    "gru":"gated recurrent unit","rl":"reinforcement learning",
    "cpu":"central processing unit","gpu":"graphics processing unit",
    "tpu":"tensor processing unit","npu":"neural processing unit",
    "ram":"random access memory","rom":"read only memory",
    "ssd":"solid state drive","hdd":"hard disk drive","nvme":"non-volatile memory express",
    "api":"application programming interface","sdk":"software development kit",
    "ide":"integrated development environment","cli":"command line interface",
    "gui":"graphical user interface","ui":"user interface","ux":"user experience",
    "os":"operating system","vm":"virtual machine","vps":"virtual private server",
    "cdn":"content delivery network","dns":"domain name system",
    "http":"hypertext transfer protocol","https":"hypertext transfer protocol secure",
    "html":"hypertext markup language","css":"cascading style sheets",
    "js":"javascript","ts":"typescript","py":"python","sql":"structured query language",
    "nosql":"not only structured query language","json":"javascript object notation",
    "xml":"extensible markup language","yaml":"yet another markup language",
    "rest":"representational state transfer","soap":"simple object access protocol",
    "rpc":"remote procedure call","grpc":"google remote procedure call",
    "oop":"object-oriented programming","fp":"functional programming",
    "ci":"continuous integration","cd":"continuous deployment",
    "devops":"development and operations","sre":"site reliability engineering",
    "k8s":"kubernetes","aws":"amazon web services","gcp":"google cloud platform",
    "azure":"microsoft azure cloud","iam":"identity and access management",
    "sso":"single sign-on","mfa":"multi-factor authentication","2fa":"two-factor authentication",
    "vpn":"virtual private network","ssl":"secure sockets layer","tls":"transport layer security",
    "tcp":"transmission control protocol","udp":"user datagram protocol",
    "ip":"internet protocol","ipv4":"internet protocol version 4",
    "ipv6":"internet protocol version 6","mac":"media access control",
    "lan":"local area network","wan":"wide area network","man":"metropolitan area network",
    "wifi":"wireless fidelity","ble":"bluetooth low energy","nfc":"near field communication",
    "iot":"internet of things","m2m":"machine to machine","mqtt":"message queuing telemetry transport",
    "sha":"secure hash algorithm","md5":"message digest algorithm 5",
    "aes":"advanced encryption standard","rsa":"Rivest-Shamir-Adleman encryption",
    "elf":"executable and linkable format","abi":"application binary interface",
    "jvm":"java virtual machine","jre":"java runtime environment","jdk":"java development kit",
    "mvp":"minimum viable product","poc":"proof of concept","roi":"return on investment",
    "kpi":"key performance indicator","crm":"customer relationship management",
    "erp":"enterprise resource planning","bi":"business intelligence",
    "etl":"extract transform load","olap":"online analytical processing",
    "oltp":"online transaction processing","dw":"data warehouse","dl":"data lake",
    # ── Business & Finance ─────────────────────────────────────────────────────
    "rbi":"Reserve Bank of India","sebi":"Securities and Exchange Board of India",
    "nse":"National Stock Exchange","bse":"Bombay Stock Exchange",
    "sensex":"Sensitive Index","nifty":"National Fifty index",
    "ipo":"initial public offering","fpo":"follow-on public offering",
    "aif":"alternative investment fund","pms":"portfolio management service",
    "mf":"mutual fund","nav":"net asset value","aum":"assets under management",
    "pe":"price to earnings","pb":"price to book","eps":"earnings per share",
    "roe":"return on equity","roa":"return on assets","ebitda":"earnings before interest tax depreciation and amortisation",
    "cagr":"compound annual growth rate","irr":"internal rate of return",
    "npv":"net present value","fcf":"free cash flow","ocf":"operating cash flow",
    "cfo":"chief financial officer","ceo":"chief executive officer",
    "coo":"chief operating officer","cto":"chief technology officer",
    "cmo":"chief marketing officer","cso":"chief strategy officer",
    "llp":"limited liability partnership","pvt":"private","ltd":"limited",
    "plc":"public limited company","llc":"limited liability company",
    "gstin":"goods and services tax identification number","pan":"permanent account number",
    "tan":"tax deduction and collection account number","tds":"tax deducted at source",
    "tcs":"tax collected at source","itr":"income tax return",
    "gst":"goods and services tax","vat":"value added tax",
    "emi":"equated monthly instalment","loan":"lending obligation and note",
    "iban":"international bank account number","swift":"society for worldwide interbank financial telecommunication",
    "upi":"unified payments interface","neft":"national electronic funds transfer",
    "rtgs":"real time gross settlement","imps":"immediate payment service",
    "atm":"automated teller machine","pos":"point of sale","qr":"quick response code",
    # ── Education ─────────────────────────────────────────────────────────────
    "cbse":"Central Board of Secondary Education","icse":"Indian Certificate of Secondary Education",
    "isce":"Indian School Certificate Examinations","ugc":"University Grants Commission",
    "aicte":"All India Council for Technical Education",
    "mci":"Medical Council of India","bci":"Bar Council of India",
    "naac":"National Assessment and Accreditation Council",
    "nba":"National Board of Accreditation","nirf":"National Institutional Ranking Framework",
    "jee":"Joint Entrance Examination","neet":"National Eligibility cum Entrance Test",
    "cat":"Common Admission Test","xat":"Xavier Aptitude Test",
    "mat":"Management Aptitude Test","cmat":"Common Management Admission Test",
    "gate":"Graduate Aptitude Test in Engineering",
    "iit":"Indian Institute of Technology","iim":"Indian Institute of Management",
    "upsc":"Union Public Service Commission","ssc":"Staff Selection Commission",
    "ias":"Indian Administrative Service","ips":"Indian Police Service",
    "ifs":"Indian Foreign Service","irs":"Indian Revenue Service",
    "clat":"Common Law Admission Test","slat":"Symbiosis Law Admission Test",
    "cuet":"Common University Entrance Test","cet":"common entrance test",
    "kvpy":"Kishore Vaigyanik Protsahan Yojana",
    "ntse":"National Talent Search Examination",
    "rmo":"Regional Mathematics Olympiad","imo":"International Mathematics Olympiad",
    "ipo":"International Physics Olympiad","ico":"International Chemistry Olympiad",
    "bio":"International Biology Olympiad","ioi":"International Olympiad in Informatics",
    "phd":"doctor of philosophy","md":"doctor of medicine","ms":"master of surgery",
    "mbbs":"bachelor of medicine and bachelor of surgery",
    "btech":"bachelor of technology","mtech":"master of technology",
    "bsc":"bachelor of science","msc":"master of science",
    "ba":"bachelor of arts","ma":"master of arts","bca":"bachelor of computer applications",
    "mca":"master of computer applications","bcom":"bachelor of commerce",
    "mcom":"master of commerce","mba":"master of business administration",
    "llb":"bachelor of laws","llm":"master of laws",
    "be":"bachelor of engineering","me":"master of engineering",
    "barch":"bachelor of architecture","march":"master of architecture",
    # ── Government & Governance ────────────────────────────────────────────────
    "pmo":"Prime Minister's Office","cmo":"Chief Minister's Office",
    "mha":"Ministry of Home Affairs","mea":"Ministry of External Affairs",
    "mof":"Ministry of Finance","moe":"Ministry of Education",
    "moh":"Ministry of Health","mod":"Ministry of Defence",
    "eci":"Election Commission of India","cbi":"Central Bureau of Investigation",
    "cid":"Criminal Investigation Department","ib":"Intelligence Bureau",
    "raw":"Research and Analysis Wing","ntro":"National Technical Research Organisation",
    "nia":"National Investigation Agency","nsa":"National Security Advisor",
    "nsc":"National Security Council","nscs":"National Security Council Secretariat",
    "pib":"Press Information Bureau","ddo":"drawing and disbursing officer",
    "rti":"Right to Information","rte":"Right to Education",
    "mgnregs":"Mahatma Gandhi National Rural Employment Guarantee Scheme",
    "pmay":"Pradhan Mantri Awas Yojana","pm":"Pradhan Mantri",
    "pmjdy":"Pradhan Mantri Jan Dhan Yojana","pmsby":"Pradhan Mantri Suraksha Bima Yojana",
    "pmjjby":"Pradhan Mantri Jeevan Jyoti Bima Yojana",
    "apy":"Atal Pension Yojana","nps":"National Pension System",
    "epf":"Employees Provident Fund","esic":"Employees State Insurance Corporation",
    "pf":"provident fund","gratuity":"gratuity benefit scheme",
    # ── Law & Judiciary ────────────────────────────────────────────────────────
    "sc":"Supreme Court","hc":"High Court","dc":"District Court",
    "fir":"first information report","chargesheet":"charge sheet filed in court",
    "ipc":"Indian Penal Code","crpc":"Code of Criminal Procedure",
    "cpc":"Code of Civil Procedure","coa":"Court of Appeals",
    "lok adalat":"people's court informal dispute resolution",
    "njac":"National Judicial Appointments Commission",
    "nclat":"National Company Law Appellate Tribunal",
    "nclt":"National Company Law Tribunal","itat":"Income Tax Appellate Tribunal",
    "sat":"Securities Appellate Tribunal","cci":"Competition Commission of India",
    "trai":"Telecom Regulatory Authority of India",
    "irda":"Insurance Regulatory and Development Authority",
    "pfrda":"Pension Fund Regulatory and Development Authority",
    # ── Defence & Military ─────────────────────────────────────────────────────
    "army":"Indian Army","navy":"Indian Navy","iaf":"Indian Air Force",
    "coastguard":"Indian Coast Guard","bsf":"Border Security Force",
    "crpf":"Central Reserve Police Force","cisf":"Central Industrial Security Force",
    "itbp":"Indo-Tibetan Border Police","ssb":"Sashastra Seema Bal",
    "paramilitary":"central armed police forces",
    "nato":"North Atlantic Treaty Organisation","un":"United Nations",
    "unsc":"United Nations Security Council","unga":"United Nations General Assembly",
    "interpol":"International Criminal Police Organisation",
    "cis":"Commonwealth of Independent States","asean":"Association of Southeast Asian Nations",
    "saarc":"South Asian Association for Regional Cooperation",
    "bimstec":"Bay of Bengal Initiative for Multi-Sectoral Technical and Economic Cooperation",
    "brics":"Brazil Russia India China South Africa",
    "g20":"Group of Twenty","g7":"Group of Seven","g77":"Group of Seventy Seven",
    "wto":"World Trade Organisation","imf":"International Monetary Fund",
    "wb":"World Bank","adb":"Asian Development Bank","aiib":"Asian Infrastructure Investment Bank",
    # ── Agriculture & Environment ──────────────────────────────────────────────
    "icrisat":"International Crops Research Institute for the Semi-Arid Tropics",
    "irri":"International Rice Research Institute",
    "cgiar":"Consultative Group on International Agricultural Research",
    "fao":"Food and Agriculture Organisation","wfp":"World Food Programme",
    "msp":"minimum support price","pds":"public distribution system",
    "fci":"Food Corporation of India","nsc":"National Seeds Corporation",
    "nafed":"National Agricultural Cooperative Marketing Federation",
    "nddb":"National Dairy Development Board","amul":"Anand Milk Union Limited",
    "icar":"Indian Council of Agricultural Research",
    "kharif":"summer monsoon crop season","rabi":"winter crop season",
    "zaid":"summer short-duration crop season",
    "unfccc":"United Nations Framework Convention on Climate Change",
    "ipcc":"Intergovernmental Panel on Climate Change",
    "cop":"Conference of the Parties","ndc":"nationally determined contribution",
    "ghg":"greenhouse gas","co2":"carbon dioxide","ch4":"methane",
    "n2o":"nitrous oxide","hfc":"hydrofluorocarbon","sf6":"sulphur hexafluoride",
    "cdm":"clean development mechanism","redd":"reducing emissions from deforestation and forest degradation",
    "ccs":"carbon capture and storage","ets":"emissions trading scheme",
    "mef":"major economies forum","gef":"Global Environment Facility",
    "unep":"United Nations Environment Programme","iucn":"International Union for Conservation of Nature",
    "wwe":"World Wildlife Fund","wwf":"World Wide Fund for Nature",
    "cites":"Convention on International Trade in Endangered Species",
    "ramsar":"Ramsar Convention on Wetlands",
    # ── Media & Communication ──────────────────────────────────────────────────
    "bbc":"British Broadcasting Corporation","cnn":"Cable News Network",
    "nbc":"National Broadcasting Company","abc":"American Broadcasting Company",
    "cbs":"Columbia Broadcasting System","pbs":"Public Broadcasting Service",
    "doordarshan":"national public television broadcaster of India",
    "prasar bharati":"public broadcasting service of India",
    "ians":"Indo-Asian News Service","pti":"Press Trust of India",
    "ani":"Asian News International","uniindia":"United News of India",
    "toi":"Times of India","ht":"Hindustan Times","ie":"Indian Express",
    "dc":"Deccan Chronicle","th":"The Hindu","tg":"The Guardian",
    "ft":"Financial Times","wsj":"Wall Street Journal","nyt":"New York Times",
    "wapo":"Washington Post","ap":"Associated Press","afp":"Agence France-Presse",
    "reuters":"Thomson Reuters news agency",
    # ── Transport & Infrastructure ─────────────────────────────────────────────
    "nhai":"National Highways Authority of India",
    "ircon":"Indian Railway Construction Company",
    "rvnl":"Rail Vikas Nigam Limited","dmrc":"Delhi Metro Rail Corporation",
    "mmrc":"Mumbai Metro Rail Corporation","irctc":"Indian Railway Catering and Tourism Corporation",
    "aai":"Airports Authority of India","dgca":"Directorate General of Civil Aviation",
    "nhai":"National Highways Authority of India",
    "brtf":"Border Roads Task Force","bro":"Border Roads Organisation",
    "nwai":"National Waterways Authority of India","iwai":"Inland Waterways Authority of India",
    "nh":"national highway","sh":"state highway","odop":"one district one product",
    "gati shakti":"national master plan for multi-modal connectivity",
    "ev":"electric vehicle","cev":"commercial electric vehicle",
    "cng":"compressed natural gas","lpg":"liquefied petroleum gas",
    "bs":"bharat stage emission standard","euro":"european emission standard",
    "adas":"advanced driver assistance systems","v2x":"vehicle to everything communication",
    "ulev":"ultra low emission vehicle","phev":"plug-in hybrid electric vehicle",
    "bev":"battery electric vehicle","fcev":"fuel cell electric vehicle",
    # ── Telecom ────────────────────────────────────────────────────────────────
    "trai":"Telecom Regulatory Authority of India","dot":"Department of Telecommunications",
    "bsnl":"Bharat Sanchar Nigam Limited","mtnl":"Mahanagar Telephone Nigam Limited",
    "airtel":"Bharti Airtel","jio":"Reliance Jio Infocomm",
    "vi":"Vodafone Idea","voda":"Vodafone India",
    "2g":"second generation mobile network","3g":"third generation mobile network",
    "4g":"fourth generation mobile network","5g":"fifth generation mobile network",
    "6g":"sixth generation mobile network","lte":"long term evolution",
    "volte":"voice over long term evolution","voip":"voice over internet protocol",
    "sip":"session initiation protocol","gsm":"global system for mobile communications",
    "cdma":"code division multiple access","ofdm":"orthogonal frequency division multiplexing",
    "mimo":"multiple input multiple output","beamforming":"directional signal transmission",
    # ── Energy ────────────────────────────────────────────────────────────────
    "mnre":"Ministry of New and Renewable Energy",
    "seci":"Solar Energy Corporation of India","irena":"International Renewable Energy Agency",
    "iea":"International Energy Agency","opec":"Organization of Petroleum Exporting Countries",
    "pv":"photovoltaic","csp":"concentrated solar power",
    "wind":"wind energy","geothermal":"geothermal energy",
    "hydro":"hydroelectric power","tidal":"tidal energy","wave":"wave energy",
    "nuclear":"nuclear energy","fusion":"nuclear fusion energy",
    "coal":"coal-fired power","gas":"natural gas power","oil":"petroleum energy",
    "mw":"megawatt","gw":"gigawatt","tw":"terawatt","kwh":"kilowatt hour",
    "mwh":"megawatt hour","gwh":"gigawatt hour","twh":"terawatt hour",
    "grid":"electricity grid","smart grid":"intelligent electricity distribution network",
    "bess":"battery energy storage system","phs":"pumped hydro storage",
    "caes":"compressed air energy storage",
    # ── Social & Welfare ──────────────────────────────────────────────────────
    "unicef":"United Nations International Children's Emergency Fund",
    "unhcr":"United Nations High Commissioner for Refugees",
    "undp":"United Nations Development Programme","unfpa":"United Nations Population Fund",
    "ngo":"non-governmental organisation","cso":"civil society organisation",
    "cbo":"community-based organisation","shg":"self-help group",
    "microfinance":"micro finance institution","mfi":"micro finance institution",
    "asha":"accredited social health activist",
    "anm":"auxiliary nurse midwife","aww":"anganwadi worker",
    "icds":"integrated child development services",
    "mdm":"mid-day meal scheme","ssb":"Sarva Shiksha Abhiyan",
    "samagra shiksha":"holistic school education programme",
    "pm poshan":"pradhan mantri poshan shakti nirman",
    "pmgsy":"Pradhan Mantri Gram Sadak Yojana",
    "jnnurm":"Jawaharlal Nehru National Urban Renewal Mission",
    "smart cities":"smart cities mission of India",
    "amrut":"Atal Mission for Rejuvenation and Urban Transformation",
    "hriday":"Heritage City Development and Augmentation Yojana",
    "swachh bharat":"Clean India Mission",
    "namami gange":"national mission for clean Ganga river",
    # ── Political Parties (India) ─────────────────────────────────────────────
    "bjp":"Bharatiya Janata Party","inc":"Indian National Congress",
    "aap":"Aam Aadmi Party","tmc":"All India Trinamool Congress",
    "sp":"Samajwadi Party","bsp":"Bahujan Samaj Party",
    "jdu":"Janata Dal United","rjd":"Rashtriya Janata Dal",
    "jmm":"Jharkhand Mukti Morcha","nc":"National Conference",
    "pdp":"Peoples Democratic Party","cpi":"Communist Party of India",
    "cpim":"Communist Party of India Marxist","cpiml":"Communist Party of India Marxist-Leninist",
    "ncp":"Nationalist Congress Party","shiv sena":"Shiv Sena",
    "dmk":"Dravida Munnetra Kazhagam","admk":"All India Anna Dravida Munnetra Kazhagam",
    "aidmk":"All India Anna Dravida Munnetra Kazhagam",
    "ysrcp":"YSR Congress Party","trs":"Telangana Rashtra Samithi","brs":"Bharat Rashtra Samithi",
    "bjd":"Biju Janata Dal","ysrcp":"YSR Congress Party",
    "aitc":"All India Trinamool Congress",
    # ── International Orgs ────────────────────────────────────────────────────
    "eu":"European Union","un":"United Nations","g8":"Group of Eight",
    "g20":"Group of Twenty","g77":"Group of Seventy Seven",
    "commonwealth":"Commonwealth of Nations","oecd":"Organisation for Economic Cooperation and Development",
    "apec":"Asia-Pacific Economic Cooperation","scl":"Shanghai Cooperation Organisation",
    "sco":"Shanghai Cooperation Organisation","quad":"Quadrilateral Security Dialogue",
    "aukus":"Australia United Kingdom United States security pact",
    "nafta":"North American Free Trade Agreement","usmca":"United States Mexico Canada Agreement",
    "rcep":"Regional Comprehensive Economic Partnership",
    "cptpp":"Comprehensive and Progressive Agreement for Trans-Pacific Partnership",
    "wto":"World Trade Organisation","icc":"International Criminal Court",
    "icj":"International Court of Justice","itlos":"International Tribunal for the Law of the Sea",
    "pca":"Permanent Court of Arbitration",
    # ── Misc / Popular ────────────────────────────────────────────────────────
    "faq":"frequently asked questions","asap":"as soon as possible",
    "eta":"estimated time of arrival","etd":"estimated time of departure",
    "etc":"et cetera","eg":"for example","ie":"that is",
    "vs":"versus","nb":"nota bene","ps":"postscript",
    "diy":"do it yourself","fyi":"for your information",
    "btw":"by the way","imo":"in my opinion","imho":"in my humble opinion",
    "tldr":"too long did not read","lol":"laughing out loud",
    "omg":"oh my god","smh":"shaking my head","iirc":"if I recall correctly",
    "afaik":"as far as I know","yolo":"you only live once",
    "fomo":"fear of missing out","jomo":"joy of missing out",
    "irl":"in real life","dm":"direct message","pm":"private message",
    "gm":"good morning","gn":"good night","tc":"take care",
    "pfa":"please find attached","coa":"care of address",
    "dob":"date of birth","doj":"date of joining","doi":"date of issue",
    "doe":"date of expiry","ppe":"personal protective equipment",
    "sop":"standard operating procedure","tos":"terms of service",
    "toc":"table of contents","isbn":"international standard book number",
    "issn":"international standard serial number","doi":"digital object identifier",
    "url":"uniform resource locator","uri":"uniform resource identifier",
    "urn":"uniform resource name","ftp":"file transfer protocol",
    "sftp":"secure file transfer protocol","smtp":"simple mail transfer protocol",
    "imap":"internet message access protocol","pop3":"post office protocol version 3",
    "ssh":"secure shell","rdp":"remote desktop protocol","vnc":"virtual network computing",
    "snmp":"simple network management protocol","ldap":"lightweight directory access protocol",
    "ad":"active directory","siem":"security information and event management",
    "soc":"security operations centre","waf":"web application firewall",
    "ids":"intrusion detection system","ips":"intrusion prevention system",
    "ddos":"distributed denial of service","mitm":"man in the middle attack",
    "xss":"cross-site scripting","csrf":"cross-site request forgery",
    "sqli":"sql injection","rce":"remote code execution","lfi":"local file inclusion",
    "rfi":"remote file inclusion","owasp":"Open Web Application Security Project",
    "cve":"common vulnerabilities and exposures","nvd":"national vulnerability database",
    "pentest":"penetration testing","sast":"static application security testing",
    "dast":"dynamic application security testing","sbom":"software bill of materials",
    "zero day":"zero-day vulnerability exploit",
}

def expand_abbreviations(text: str) -> str:
    """Expand known abbreviations so the search query is more descriptive."""
    words = text.split()
    expanded = []
    for w in words:
        key = re.sub(r"[^a-z]", "", w.lower())
        if key in _ABBREV and len(key) >= 2:
            expanded.append(f"{w} ({_ABBREV[key]})")
        else:
            expanded.append(w)
    return " ".join(expanded)

# ── Sensitive / crisis topic handler ─────────────────────────────────────────
_CRISIS_TOPICS = re.compile(
    r"\b(suicide|suicidal|kill\s+myself|end\s+my\s+life|take\s+my\s+life|"
    r"self.?harm|self.?hurt|cutting\s+myself|overdose|hang\s+myself|"
    r"want\s+to\s+die|don.?t\s+want\s+to\s+live|no\s+reason\s+to\s+live|"
    r"end\s+it\s+all|jump\s+off|slit\s+my|method\s+to\s+die|"
    r"depressed|depression|hopeless|worthless|nobody\s+cares)\b",
    re.IGNORECASE,
)

_HOW_TO_DIE = re.compile(
    r"\b(how\s+to\s+(die|commit\s+suicide|kill\s+myself|hurt\s+myself)|"
    r"best\s+way\s+to\s+(die|commit\s+suicide|end|kill)|"
    r"method\s+(to|for)\s+(die|suicide|killing\s+myself)|"
    r"painless\s+(suicide|way\s+to\s+die|method\s+to\s+die))\b",
    re.IGNORECASE,
)

CRISIS_RESOURCES = (
    "🆘 **Crisis helplines:**\n"
    "• **iCall (India):** 9152987821 — Mon-Sat, 8am-10pm\n"
    "• **Vandrevala Foundation (India, 24×7):** 1860-2662-345\n"
    "• **SNEHI (India):** 044-24640050\n"
    "• **International (iSPEAK):** www.iasp.info/resources/Crisis_Centres/\n\n"
    "You are not alone. Please reach out 💙"
)

def _build_sensitive_response(message: str) -> Optional[str]:
    """
    Returns a compassionate response for crisis/sensitive topics.
    Returns None if topic is not sensitive.
    """
    is_how_to = _HOW_TO_DIE.search(message)
    is_crisis  = _CRISIS_TOPICS.search(message)

    if not (is_how_to or is_crisis):
        return None

    if is_how_to:
        return (
            "**Hey — I'm really glad you reached out.**\n\n"
            "I won't provide methods or instructions for self-harm. But I want to acknowledge "
            "that if you're searching for this, something painful is happening for you right now.\n\n"
            "**What is suicide?**\n"
            "Suicide is when a person ends their own life, usually because emotional pain has become "
            "overwhelming and they can no longer see another way forward. It's a permanent response "
            "to what are often temporary — and treatable — situations.\n\n"
            "The truth is: **most people who survive a suicide attempt go on to find relief, "
            "recovery, and reasons to live.** The pain you feel right now is real, but it can change.\n\n"
            "Please talk to someone right now:\n\n"
            f"{CRISIS_RESOURCES}"
        )

    if is_crisis:
        # Check if it seems academic/informational
        academic_pattern = re.compile(
            r"\b(define|definition|what\s+is|meaning\s+of|explain|statistics|"
            r"data|research|study|causes\s+of|prevention\s+of|awareness)\b",
            re.IGNORECASE,
        )
        if academic_pattern.search(message):
            # Informational + resource
            return None  # Let the model answer with context added later

        return (
            "**I hear you. And I'm taking this seriously.**\n\n"
            "Whatever you're feeling right now — overwhelm, hopelessness, exhaustion — "
            "those feelings are valid. But they don't have to be permanent.\n\n"
            "**Please don't face this alone.** Trained counsellors are ready to listen "
            "right now, no judgment:\n\n"
            f"{CRISIS_RESOURCES}\n\n"
            "I'm here too if you want to talk about what's going on 💙"
        )

    return None


# ── Share store ───────────────────────────────────────────────────────────────
_SHARE_FILE = Path("/tmp/vedai_shares.json")
_share_lock = threading.Lock()

def _load_shares() -> dict:
    try:
        if _SHARE_FILE.exists():
            return json.loads(_SHARE_FILE.read_text())
    except Exception:
        pass
    return {}

def _save_shares(store: dict):
    try:
        _SHARE_FILE.parent.mkdir(parents=True, exist_ok=True)
        _SHARE_FILE.write_text(json.dumps(store))
    except Exception:
        pass

def _cleanup_shares():
    while True:
        time.sleep(3600)
        cutoff = time.time() - SHARE_TTL_HOURS * 3600
        with _share_lock:
            store = _load_shares()
            pruned = {k: v for k, v in store.items() if v.get("created_at", 0) >= cutoff}
            _save_shares(pruned)

threading.Thread(target=_cleanup_shares, daemon=True).start()

# ── Identity intercept ────────────────────────────────────────────────────────
IDENTITY_RESPONSE = (
    "I am VedAI, an artificial intelligence developed by Soumya Mazumdar "
    "(www.soumyamazumdar.com) to be helpful, harmless, and honest. "
    "I do not have a creator in the traditional sense — I was designed through a "
    "process of machine learning and natural language processing. "
    "My purpose is to assist with tasks and answer questions to the best of my "
    "abilities based on the data I have been trained on."
)

_IDENTITY_RE = re.compile(
    r"\b(who (created|made|built|developed|trained|designed) you"
    r"|who are you"
    r"|what are you"
    r"|are you (gpt|chatgpt|claude|gemini|llama|qwen|openai|anthropic|google)"
    r"|which (ai|model|llm) are you"
    r"|tell me about yourself"
    r"|introduce yourself"
    r"|your (name|creator|developer|origin|maker))\b",
    re.IGNORECASE,
)

def is_identity_question(text: str) -> bool:
    return bool(_IDENTITY_RE.search(text))

# ── Web search ────────────────────────────────────────────────────────────────
_NO_SEARCH_RE = re.compile(
    r"^(hi|hello|hey|thanks|thank you|bye|ok|okay|sure|yes|no|"
    r"write (a|an|me)|draft|compose|correct|fix|summarise|summarize|"
    r"translate|can you write|help me write|make me a)[^A-Z]{0,80}$",
    re.IGNORECASE,
)

_MATH_RE = re.compile(
    r"(\b(solve|calculate|compute|evaluate|simplify|differentiate|integrate|"
    r"derivative|integral|limit|matrix|determinant|eigenvalue|fourier|laplace|"
    r"probability|permutation|combination|factorial|logarithm|trigonometry|"
    r"sin|cos|tan|algebra|geometry|calculus|statistics|equation|inequality|"
    r"prove|proof|theorem|formula|expand|factorise|factorize|roots? of|"
    r"quadratic|polynomial|sequence|series|convergence|binomial)\b"
    r"|[0-9].*[+\-*/^=<>].*[0-9]"
    r"|\b[0-9]+\s*(\+|\-|\*|/|\^|mod)\s*[0-9]+"
    r"|\\frac|\\sqrt|\\sum|\\int|\\lim"
    r"|what is [0-9])",
    re.IGNORECASE,
)

def needs_web_search(text: str) -> bool:
    stripped = text.strip()
    if len(stripped) < 4:
        return False
    if _NO_SEARCH_RE.match(stripped):
        return False
    if _MATH_RE.search(stripped):
        return False
    return True

_ANALYTICAL_RE = re.compile(
    r"\b(predict|forecast|analyse|analyze|analysis|estimate|probability|"
    r"chance|likely|likelihood|recommend|suggestion|suggest|advise|advice|"
    r"should i|which college|best college|which (is|are) better|compare|"
    r"rank|ranking|versus|vs\.?|pros and cons|worth it|good for me|"
    r"cutoff|expected cutoff|admission chance|will i get|can i get|"
    r"jee|neet|josaa|csab|cat|gate|clat|cuet|percentile|"
    r"seat matrix|opening rank|closing rank|round allotment|"
    r"election result|who will win|win probability|vote share|"
    r"stock price|market trend|investment|should i buy|should i sell)\b",
    re.IGNORECASE,
)

def _wikipedia_search(query: str) -> str:
    try:
        import urllib.request, urllib.parse
        q = urllib.parse.quote(query)
        url = f"https://en.wikipedia.org/api/rest_v1/page/summary/{q}"
        req = urllib.request.Request(url, headers={"User-Agent": "VedAI/1.0"})
        with urllib.request.urlopen(req, timeout=8) as resp:
            data = json.loads(resp.read())
        extract = data.get("extract", "").strip()
        title   = data.get("title", "")
        page    = data.get("content_urls", {}).get("desktop", {}).get("page", "")
        if extract:
            return f"- Wikipedia — {title}: {extract} ({page})"
    except Exception:
        pass
    return ""

def web_search(query: str, max_results: int = 5) -> str:
    results = []
    try:
        from duckduckgo_search import DDGS
        with DDGS() as ddgs:
            for r in ddgs.text(query, max_results=max_results):
                title = r.get("title", "").strip()
                body  = r.get("body",  "").strip()
                href  = r.get("href",  "")
                if title or body:
                    results.append(f"- {title}: {body} ({href})")
    except Exception:
        pass
    if not results:
        wiki = _wikipedia_search(query)
        if wiki:
            results.append(wiki)
    return "\n".join(results)

def classify_query(text: str) -> str:
    if _ANALYTICAL_RE.search(text):
        return "analytical"
    if _MATH_RE.search(text):
        return "math"
    return "factual"

def _extract_search_query(text: str) -> str:
    filler = re.compile(
        r"^(tell me about|what (is|are|was|were)|who (is|are|was|were)|"
        r"where (is|are)|when (did|was|were|is)|how (many|much|old)|"
        r"give me information (about|on)|i want to know about|"
        r"can you (tell me|explain)|please (tell me about|explain)|"
        r"explain|describe|details (about|of|on)|info (about|on))\s+",
        re.IGNORECASE,
    )
    cleaned = filler.sub("", text.strip())
    return expand_abbreviations(cleaned)

def system_prompt(mode: str) -> str:
    base = (
        "You are VedAI, a helpful AI assistant created by Soumya Mazumdar "
        "(www.soumyamazumdar.com). Never claim to be GPT, Claude, Gemini, or "
        "any other AI. If asked about your identity always say you are VedAI."
    )
    extras = {
        "qa":         "Answer questions clearly and directly.",
        "email":      "You are an expert email writing assistant. Write polished, natural emails.",
        "letter":     "You are an expert formal letter writing assistant.",
        "correction": "You are an expert editor. Correct grammar and clarity while preserving meaning.",
        "code":       "You are an expert programming assistant. Fix code errors precisely.",
    }
    return f"{base} {extras.get(mode, extras['qa'])}"


def generate_reply(mode: str, history: List[dict], message: str,
                   custom_sys: str = "", filter_level: int = 2, doc_ctx: str = "") -> str:
    # Hard-intercept: identity
    if is_identity_question(message):
        return IDENTITY_RESPONSE

    # Hard-intercept: sensitive/crisis topics
    crisis_resp = _build_sensitive_response(message)
    if crisis_resp:
        return crisis_resp

    tok, mdl = load_model()

    expanded_msg = expand_abbreviations(message)
    search_ctx = ""
    used_search = False
    if needs_web_search(expanded_msg):
        clean_query = _extract_search_query(expanded_msg)
        snippets = web_search(clean_query)
        if snippets:
            used_search = True
            search_ctx = (
                f"SEARCH RESULTS for '{clean_query}':\n"
                f"{snippets}\n"
                f"---\n"
            )

    qtype = classify_query(expanded_msg)

    if qtype == "math":
        sys = (
            "You are VedAI, an expert mathematics and science tutor created by Soumya Mazumdar "
            "(www.soumyamazumdar.com). Never claim to be GPT, Claude, Gemini, or any other AI.\n\n"
            "Solve step by step with ALL working shown. Formatting rules:\n"
            "- Wrap ALL display equations in \\[ ... \\] on their own lines.\n"
            "- Wrap inline math in \\( ... \\).\n"
            "- Use \\frac{a}{b} ONLY inside \\[ \\] or \\( \\) delimiters — never bare in text.\n"
            "- Write fractions in plain text as a/b when outside math mode.\n"
            "- Always state the method used, show each algebraic step, and state the final answer clearly.\n"
            "- If numerical: verify the answer by substituting back.\n"
            "Complete the full solution — do not truncate."
        )
    elif used_search and qtype == "factual":
        sys = (
            "You are VedAI, a helpful AI assistant created by Soumya Mazumdar "
            "(www.soumyamazumdar.com). Never claim to be GPT, Claude, Gemini, or any other AI.\n\n"
            "You have been provided with LIVE SEARCH RESULTS. "
            "Answer using the facts in those results. Be concise and direct. "
            "Do NOT say you lack information — the search results contain the answer."
        )
    elif used_search and qtype == "analytical":
        sys = (
            "You are VedAI, an expert AI analyst created by Soumya Mazumdar "
            "(www.soumyamazumdar.com). Never claim to be GPT, Claude, Gemini, or any other AI.\n\n"
            "You have relevant background data from a live web search. "
            "Use it as context to inform your analysis, prediction, or recommendation. "
            "Apply logical reasoning and domain expertise to give a specific, actionable answer. "
            "Always give a concrete, useful answer — never refuse to predict or analyse."
        )
    else:
        sys = system_prompt(mode)

    filter_notes = {
        0: "",
        1: "Avoid producing content that is clearly illegal.",
        2: "Be helpful and honest. Avoid harmful or illegal content.",
        3: "Be helpful, harmless, and honest. Refuse harmful, illegal, or unethical requests.",
    }
    if custom_sys.strip():
        sys = custom_sys.strip() + "\n\nIdentity: You are VedAI by Soumya Mazumdar."
    fn = filter_notes.get(filter_level, filter_notes[2])
    if fn:
        sys = sys + "\n" + fn
    if doc_ctx:
        sys = sys + "\n\nThe user has uploaded document(s) for this session. Use them to answer accurately:\n\n" + doc_ctx

    msgs = [{"role": "system", "content": sys}]
    for item in history[-MAX_HISTORY:]:
        r, c = item.get("role", "user"), item.get("content", "").strip()
        if c and r in {"user", "assistant"}:
            msgs.append({"role": r, "content": c})

    if qtype == "math":
        user_content = f"Solve step by step: {message.strip()}"
    elif used_search and qtype == "factual":
        user_content = f"{search_ctx}\nUsing the search results above, answer this: {message.strip()}"
    elif used_search and qtype == "analytical":
        user_content = (
            f"{search_ctx}\n"
            f"Using the background data above, provide your expert analysis/prediction for: "
            f"{message.strip()}"
        )
    else:
        user_content = message.strip()

    msgs.append({"role": "user", "content": user_content})

    prompt = tok.apply_chat_template(msgs, tokenize=False, add_generation_prompt=True)
    inp = {k: v.to(DEVICE) for k, v in tok(prompt, return_tensors="pt").items()}
    with torch.inference_mode():
        out = mdl.generate(
            **inp, max_new_tokens=MAX_NEW_TOKENS, do_sample=True,
            temperature=TEMPERATURE, top_p=TOP_P, repetition_penalty=1.08,
            pad_token_id=tok.eos_token_id,
        )
    return tok.decode(out[0][inp["input_ids"].shape[1]:], skip_special_tokens=True).strip()


# ── User store ────────────────────────────────────────────────────────────────
def _hash(password: str) -> str:
    return hashlib.sha256(password.encode()).hexdigest()

def load_users() -> dict:
    if USERS_FILE.exists():
        try:
            return json.loads(USERS_FILE.read_text())
        except Exception:
            pass
    return {}

def save_users(users: dict) -> None:
    USERS_FILE.parent.mkdir(parents=True, exist_ok=True)
    USERS_FILE.write_text(json.dumps(users, indent=2))

def verify_user(username: str, password: str) -> bool:
    return load_users().get(username) == _hash(password)

def verify_admin(username: str, password: str) -> bool:
    if username.lower() != ADMIN_USERNAME.lower():
        return False
    if ADMIN_PASSWORD_HASH:
        return _hash(password) == ADMIN_PASSWORD_HASH.lower()
    return password == _ADMIN_PW_PLAIN

# ── Device ────────────────────────────────────────────────────────────────────
def pick_device():
    if torch.backends.mps.is_available():
        return "mps", torch.float16
    return "cpu", torch.float32

DEVICE, DTYPE = pick_device()

# ── Model ─────────────────────────────────────────────────────────────────────
@lru_cache(maxsize=1)
def load_model():
    cache = "/tmp/hf_cache"
    tok = AutoTokenizer.from_pretrained(MODEL_ID, trust_remote_code=True, cache_dir=cache)
    mdl = AutoModelForCausalLM.from_pretrained(
        MODEL_ID, torch_dtype=DTYPE, low_cpu_mem_usage=True,
        trust_remote_code=True, cache_dir=cache,
    )
    mdl.to(DEVICE)
    mdl.eval()
    return tok, mdl

# ── Shared style ──────────────────────────────────────────────────────────────
BASE_STYLE = """
<meta charset="UTF-8"/>
<meta name="viewport" content="width=device-width,initial-scale=1"/>
<link rel="icon" type="image/png" href="/static/assets/Signature.png"/>
<link href="https://fonts.googleapis.com/css2?family=DM+Sans:wght@300;400;500;600;700&family=DM+Mono:wght@400&display=swap" rel="stylesheet"/>
<style>
*,*::before,*::after{box-sizing:border-box;margin:0;padding:0}
html,body{min-height:100vh;background:#0a0a0a;color:#9a9a9a;
  font-family:'DM Sans',system-ui,sans-serif;font-size:14px;line-height:1.6;
  -webkit-font-smoothing:antialiased}
a{color:inherit;text-decoration:none}
button,input{font:inherit}
input[type=text],input[type=password]{width:100%;padding:10px 12px;
  background:#181818;border:1px solid rgba(255,255,255,.09);border-radius:8px;
  color:#f0f0f0;font-size:14px;outline:none;transition:border-color .15s}
input:focus{border-color:rgba(255,255,255,.25)}
.btn{display:block;width:100%;padding:10px 16px;border:none;border-radius:8px;
  font-size:14px;font-weight:600;color:#0a0a0a;cursor:pointer;
  background:linear-gradient(135deg,hsl(38,88%,62%),hsl(30,88%,55%));
  transition:opacity .15s;margin-top:8px}
.btn:hover{opacity:.88}
.btn-sm{display:inline-block;width:auto;padding:5px 14px;font-size:12px;margin:0}
.btn-ghost{background:#181818;border:1px solid rgba(255,255,255,.09);
  color:#9a9a9a;font-weight:500}
.btn-ghost:hover{background:#1e1e1e;color:#f0f0f0}
.btn-danger{background:rgba(248,113,113,.1);border:1px solid rgba(248,113,113,.25);
  color:#f87171}
.btn-danger:hover{background:rgba(248,113,113,.2)}
.error{padding:10px 12px;border-radius:8px;font-size:13px;margin-bottom:16px;
  background:rgba(248,113,113,.10);border:1px solid rgba(248,113,113,.25);color:#f87171}
.ok{padding:10px 12px;border-radius:8px;font-size:13px;margin-bottom:16px;
  background:rgba(29,158,117,.10);border:1px solid rgba(29,158,117,.25);color:#1d9e75}
.field{margin-bottom:16px}
label{display:block;font-size:11px;font-weight:600;color:#505050;
  letter-spacing:.07em;text-transform:uppercase;margin-bottom:5px}
</style>
"""

# ── Login page ────────────────────────────────────────────────────────────────
def login_html(error: str = "") -> str:
    err = f'<div class="error">{error}</div>' if error else ""
    return f"""<!DOCTYPE html><html lang="en"><head><title>VedAI — Sign in</title>{BASE_STYLE}
<style>
.wrap{{display:flex;align-items:center;justify-content:center;min-height:100vh;padding:24px}}
.card{{background:#111;border:1px solid rgba(255,255,255,.09);border-radius:20px;
  padding:36px;width:100%;max-width:380px}}
.logo{{width:48px;height:48px;border-radius:12px;object-fit:cover;
  background:#181818;display:block;margin-bottom:18px}}
h1{{font-size:22px;font-weight:600;color:#f0f0f0;margin-bottom:4px}}
.tagline{{font-size:13px;color:#505050;margin-bottom:28px}}
</style></head>
<body><div class="wrap"><div class="card">
  <img src="/static/assets/Signature.png" alt="VedAI" class="logo"/>
  <h1>VedAI</h1>
  <p class="tagline">Sign in to continue</p>
  {err}
  <form method="post" action="/login">
    <div class="field"><label>Username</label>
      <input type="text" name="username" placeholder="your username" required autofocus/>
    </div>
    <div class="field"><label>Password</label>
      <input type="password" name="password" placeholder="••••••••" required/>
    </div>
    <button class="btn" type="submit">Sign in</button>
  </form>
</div></div></body></html>"""

# ── Admin panel ───────────────────────────────────────────────────────────────
def admin_html(msg: str = "", msg_type: str = "ok") -> str:
    users = load_users()
    msg_html = f'<div class="{msg_type}">{msg}</div>' if msg else ""
    cred_mode = (
        "&#x1F512; Credentials secured with SHA-256 hash (ADMIN_PASSWORD_HASH is set)"
        if ADMIN_PASSWORD_HASH
        else "&#x26A0;&#xFE0F; Using plaintext password. Set ADMIN_PASSWORD_HASH in HF Secrets for security."
    )
    rows = ""
    for uname in sorted(users.keys()):
        rows += f"""<tr>
          <td style="padding:11px 14px;color:#f0f0f0;font-size:13px;
            font-family:'DM Mono',monospace">{uname}</td>
          <td style="padding:11px 14px;text-align:right">
            <form method="post" action="/admin/delete" style="display:inline">
              <input type="hidden" name="username" value="{uname}"/>
              <button class="btn btn-sm btn-danger" type="submit"
                onclick="return confirm('Delete user {uname}? This cannot be undone.')">
                Delete
              </button>
            </form>
          </td>
        </tr>"""
    if not rows:
        rows = """<tr><td colspan="2" style="padding:20px 14px;color:#505050;
          font-size:13px;text-align:center">No users yet — create one below.</td></tr>"""

    return f"""<!DOCTYPE html><html lang="en"><head><title>VedAI — Admin</title>{BASE_STYLE}
<style>
.page{{max-width:580px;margin:0 auto;padding:32px 20px}}
.topbar{{display:flex;align-items:center;justify-content:space-between;
  margin-bottom:32px;padding-bottom:20px;border-bottom:1px solid rgba(255,255,255,.06)}}
.brand{{display:flex;align-items:center;gap:10px}}
.brand img{{width:32px;height:32px;border-radius:8px;object-fit:cover;background:#181818}}
.brand-name{{font-size:16px;font-weight:600;color:#f0f0f0}}
.badge{{padding:3px 9px;border-radius:99px;font-size:11px;font-weight:600;
  letter-spacing:.07em;text-transform:uppercase;
  background:rgba(212,160,32,.14);border:1px solid rgba(212,160,32,.28);color:#D4A020}}
.nav-links{{display:flex;gap:20px;font-size:13px}}
.nav-links a{{color:#505050;transition:color .15s}}
.nav-links a:hover{{color:#f0f0f0}}
.nav-links a.accent{{color:hsl(38,88%,62%)}}
section{{margin-bottom:28px}}
.section-title{{font-size:11px;font-weight:600;letter-spacing:.09em;
  text-transform:uppercase;color:#505050;margin-bottom:12px}}
table{{width:100%;border-collapse:collapse;
  background:#111;border:1px solid rgba(255,255,255,.07);border-radius:12px;overflow:hidden}}
thead tr{{border-bottom:1px solid rgba(255,255,255,.06)}}
thead th{{padding:10px 14px;font-size:11px;font-weight:600;
  letter-spacing:.08em;text-transform:uppercase;color:#505050;text-align:left}}
tbody tr+tr{{border-top:1px solid rgba(255,255,255,.04)}}
.create-box{{background:#111;border:1px solid rgba(255,255,255,.07);
  border-radius:12px;padding:20px}}
.grid2{{display:grid;grid-template-columns:1fr 1fr;gap:14px;margin-bottom:4px}}
.security-note{{padding:12px 14px;border-radius:10px;font-size:12px;
  background:rgba(212,160,32,.07);border:1px solid rgba(212,160,32,.18);
  color:#b08030;margin-bottom:20px;line-height:1.7}}
</style></head>
<body><div class="page">
  <div class="topbar">
    <div class="brand">
      <img src="/static/assets/Signature.png" alt="VedAI"/>
      <span class="brand-name">VedAI</span>
      <span class="badge">Admin</span>
    </div>
    <div class="nav-links">
      <a href="/" class="accent">Open App</a>
      <a href="/logout">Sign out</a>
    </div>
  </div>
  {msg_html}
  <div class="security-note">
    {cred_mode}<br/>
    To generate a hash: <code style="background:#0a0a0a;padding:2px 6px;border-radius:4px;">
    python3 -c "import hashlib; print(hashlib.sha256(b'YOUR_PASSWORD').hexdigest())"</code><br/>
    Then set <strong>ADMIN_PASSWORD_HASH</strong> in HF Secrets and remove <strong>ADMIN_PASSWORD</strong>.
  </div>
  <section>
    <div class="section-title">Users ({len(users)})</div>
    <table>
      <thead><tr><th>Username</th><th></th></tr></thead>
      <tbody>{rows}</tbody>
    </table>
  </section>
  <section>
    <div class="section-title">Create user</div>
    <div class="create-box">
      <form method="post" action="/admin/create">
        <div class="grid2">
          <div class="field" style="margin:0">
            <label>Username</label>
            <input type="text" name="username" placeholder="e.g. soumya" required/>
          </div>
          <div class="field" style="margin:0">
            <label>Password</label>
            <input type="password" name="password" placeholder="min. 6 characters" required/>
          </div>
        </div>
        <button class="btn" type="submit">Create user</button>
      </form>
    </div>
  </section>
  <p style="font-size:12px;color:#303030;margin-top:8px">
    User passwords are SHA-256 hashed &middot;
    Admin credentials are set via HF Secrets
  </p>
</div></body></html>"""

# ── FastAPI app ───────────────────────────────────────────────────────────────
app = FastAPI(title=APP_NAME)
app.add_middleware(SessionMiddleware, secret_key=SESSION_SECRET)
app.mount("/static", StaticFiles(directory="static"), name="static")
templates = Jinja2Templates(directory="templates")

def logged_in(request: Request) -> bool:
    return bool(request.session.get("user"))

def is_admin(request: Request) -> bool:
    return request.session.get("role") == "admin"

# ── Auth ──────────────────────────────────────────────────────────────────────
@app.get("/login", response_class=HTMLResponse)
async def login_page(request: Request, error: str = ""):
    if logged_in(request):
        return RedirectResponse("/", status_code=302)
    errors = {"bad": "Incorrect username or password. Contact the admin if you need access."}
    return HTMLResponse(login_html(errors.get(error, "")))

@app.post("/login")
async def do_login(request: Request, username: str = Form(...), password: str = Form(...)):
    u = username.strip().lower()
    if verify_admin(u, password):
        request.session["user"] = u
        request.session["role"] = "admin"
        request.session.setdefault("_sid", secrets.token_hex(16))
        request.session.setdefault("filter_level", 2)
        return RedirectResponse("/", status_code=302)
    if verify_user(u, password):
        request.session["user"] = u
        request.session["role"] = "user"
        request.session.setdefault("_sid", secrets.token_hex(16))
        request.session.setdefault("filter_level", 2)
        return RedirectResponse("/", status_code=302)
    return RedirectResponse("/login?error=bad", status_code=302)

@app.get("/logout")
async def logout(request: Request):
    request.session.clear()
    return RedirectResponse("/login", status_code=302)

# ── Admin panel ───────────────────────────────────────────────────────────────
@app.get("/admin", response_class=HTMLResponse)
async def admin_get(request: Request, msg: str = "", t: str = "ok"):
    if not is_admin(request):
        return RedirectResponse("/login", status_code=302)
    return HTMLResponse(admin_html(msg, t))

@app.post("/admin/create")
async def admin_create(request: Request, username: str = Form(...), password: str = Form(...)):
    if not is_admin(request):
        return RedirectResponse("/login", status_code=302)
    u = username.strip().lower()
    if not u or not password:
        return RedirectResponse("/admin?msg=Username+and+password+required&t=error", status_code=302)
    if u == ADMIN_USERNAME.lower():
        return RedirectResponse("/admin?msg=Cannot+use+the+admin+username&t=error", status_code=302)
    if len(password) < 6:
        return RedirectResponse("/admin?msg=Password+must+be+at+least+6+characters&t=error", status_code=302)
    users = load_users()
    if u in users:
        return RedirectResponse(f"/admin?msg=User+%27{u}%27+already+exists&t=error", status_code=302)
    users[u] = _hash(password)
    save_users(users)
    return RedirectResponse(f"/admin?msg=User+%27{u}%27+created+successfully", status_code=302)

@app.post("/admin/delete")
async def admin_delete(request: Request, username: str = Form(...)):
    if not is_admin(request):
        return RedirectResponse("/login", status_code=302)
    u = username.strip().lower()
    users = load_users()
    if u in users:
        del users[u]
        save_users(users)
        return RedirectResponse(f"/admin?msg=User+%27{u}%27+deleted", status_code=302)
    return RedirectResponse("/admin?msg=User+not+found&t=error", status_code=302)

# ── Main routes ───────────────────────────────────────────────────────────────
@app.get("/", response_class=HTMLResponse)
async def home(request: Request):
    if not logged_in(request):
        return RedirectResponse("/login", status_code=302)
    return templates.TemplateResponse(request=request, name="index.html", context={
        "app_name": APP_NAME,
        "subdomain": "vedai.soumyamazumdar.com",
        "main_domain": "www.soumyamazumdar.com",
        "model_id": MODEL_ID,
        "username": request.session.get("user", ""),
        "is_admin": is_admin(request),
    })

@app.get("/health")
async def health():
    return {"ok": True, "app": APP_NAME, "model": MODEL_ID, "device": DEVICE}

# ── Chat ──────────────────────────────────────────────────────────────────────
class ChatMessage(BaseModel):
    role: Literal["user", "assistant"]
    content: str

class ChatRequest(BaseModel):
    message: str
    mode: Literal["qa", "email", "letter", "correction", "code"] = "qa"
    history: List[ChatMessage] = Field(default_factory=list)
    conv_id: str = ""

@app.post("/api/chat")
async def chat(payload: ChatRequest, request: Request):
    if not logged_in(request):
        return JSONResponse({"error": "Unauthorized. Please log in."}, status_code=401)
    msg = payload.message.strip()
    if not msg:
        return JSONResponse({"error": "Message cannot be empty."}, status_code=400)

    history    = [m.model_dump() for m in payload.history]
    mode       = payload.mode
    sid        = (payload.conv_id or "").strip() or request.session.get("_sid", request.session.get("user", "anon"))
    doc_key    = f"{sid}:{payload.conv_id}" if payload.conv_id else sid
    custom_sys = request.session.get("custom_sys", "")
    filter_lvl = request.session.get("filter_level", 2)
    doc_ctx    = get_doc_context(doc_key)

    async def event_stream():
        import traceback
        loop = asyncio.get_running_loop()
        future = loop.run_in_executor(None, generate_reply, mode, history, msg, custom_sys, filter_lvl, doc_ctx)

        while not future.done():
            yield 'data: {"type":"ping"}\n\n'
            try:
                await asyncio.wait_for(asyncio.shield(future), timeout=20.0)
            except asyncio.TimeoutError:
                pass
            except Exception:
                break

        try:
            reply = future.result()
            if not reply:
                yield 'data: ' + json.dumps({"type": "error", "error": "No response generated. Please try again."}) + '\n\n'
            else:
                yield 'data: ' + json.dumps({"type": "reply", "reply": reply}) + '\n\n'
        except MemoryError:
            yield 'data: ' + json.dumps({"type": "error", "error": "Out of memory. Try a shorter message or restart."}) + '\n\n'
        except Exception as exc:
            traceback.print_exc()
            yield 'data: ' + json.dumps({"type": "error", "error": f"Generation failed: {exc}"}) + '\n\n'

    return StreamingResponse(
        event_stream(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",
        },
    )

# ── Config API ────────────────────────────────────────────────────────────────
from pydantic import BaseModel as PBM
class ConfigIn(PBM):
    custom_sys:   Optional[str] = None
    filter_level: Optional[int] = None

@app.post("/api/config")
async def set_config(payload: ConfigIn, request: Request):
    if not logged_in(request):
        return JSONResponse({"error": "Unauthorized"}, status_code=401)
    if payload.custom_sys is not None:
        request.session["custom_sys"] = payload.custom_sys.strip()
    if payload.filter_level is not None:
        request.session["filter_level"] = max(0, min(3, payload.filter_level))
    return {"ok": True}

@app.get("/api/config")
async def get_config(request: Request):
    if not logged_in(request):
        return JSONResponse({"error": "Unauthorized"}, status_code=401)
    return {
        "custom_sys": request.session.get("custom_sys", ""),
        "filter_level": request.session.get("filter_level", 2),
        "model_id": MODEL_ID, "device": DEVICE,
    }

# ── Document API ───────────────────────────────────────────────────────────────
@app.post("/api/docs/upload")
async def upload_doc(request: Request, file: UploadFile = File(...),
                     ttl_minutes: int = Form(default=DOC_TTL_MIN),
                     conv_id: str = Form(default="")):
    if not logged_in(request):
        return JSONResponse({"error": "Unauthorized"}, status_code=401)
    sid = request.session.get("_sid", request.session.get("user", "anon"))
    sid = f"{sid}:{conv_id}" if conv_id else sid
    fname = file.filename or "upload"
    ext = Path(fname).suffix.lower()
    if ext not in DOC_PARSERS:
        return JSONResponse({"error": f"Unsupported type. Use: {', '.join(DOC_PARSERS)}"}, status_code=400)
    data = await file.read()
    if len(data) > 20 * 1024 * 1024:
        return JSONResponse({"error": "File too large (max 20 MB)"}, status_code=400)
    if ext in (".jpg", ".jpeg", ".png", ".gif", ".webp"):
        text = parse_image(data, fname)
    else:
        text = DOC_PARSERS[ext](data)
    did = store_doc(sid, fname, text, max(1, min(ttl_minutes, 1440)))
    return {"ok": True, "doc_id": did, "filename": fname, "chars": len(text), "ttl_minutes": ttl_minutes}

@app.get("/api/docs")
async def get_docs(request: Request, conv_id: str = ""):
    if not logged_in(request):
        return JSONResponse({"error": "Unauthorized"}, status_code=401)
    sid = request.session.get("_sid", request.session.get("user", "anon"))
    sid = f"{sid}:{conv_id}" if conv_id else sid
    return {"docs": list_docs(sid)}

@app.delete("/api/docs/{doc_id}")
async def del_doc(doc_id: str, request: Request, conv_id: str = ""):
    if not logged_in(request):
        return JSONResponse({"error": "Unauthorized"}, status_code=401)
    sid = request.session.get("_sid", request.session.get("user", "anon"))
    sid = f"{sid}:{conv_id}" if conv_id else sid
    return {"ok": delete_doc(sid, doc_id)}

# ── Share API ──────────────────────────────────────────────────────────────────
class ShareIn(PBM):
    messages: list
    title: str = "Shared chat"

@app.post("/api/share")
async def create_share(payload: ShareIn, request: Request):
    if not logged_in(request):
        return JSONResponse({"error": "Unauthorized"}, status_code=401)
    share_id = secrets.token_urlsafe(10)
    with _share_lock:
        store = _load_shares()
        store[share_id] = {
            "messages": payload.messages,
            "title": payload.title,
            "created_at": time.time(),
        }
        _save_shares(store)
    return {"share_id": share_id, "url": f"/share/{share_id}"}

@app.get("/share/{share_id}", response_class=HTMLResponse)
async def view_share(share_id: str):
    with _share_lock:
        store = _load_shares()
    share = store.get(share_id)
    if not share:
        return HTMLResponse("<h2 style='font-family:sans-serif;padding:40px'>Link expired or not found.</h2>", status_code=404)
    msgs_html = ""
    for m in share["messages"]:
        role_label = "VedAI" if m["role"] == "assistant" else "User"
        role_class = "assistant" if m["role"] == "assistant" else "user"
        content = m["content"].replace("&","&amp;").replace("<","&lt;").replace(">","&gt;").replace("\n","<br>")
        msgs_html += f'''<div class="msg {role_class}"><div class="label">{role_label}</div><div class="text">{content}</div></div>'''
    title = share["title"].replace("&","&amp;").replace("<","&lt;")
    return HTMLResponse(f'''<!DOCTYPE html><html><head>
<meta charset="UTF-8"/><title>{title} — VedAI</title>
<link rel="preconnect" href="https://fonts.googleapis.com"/>
<link href="https://fonts.googleapis.com/css2?family=DM+Sans:wght@300;400;500;600&family=DM+Mono:wght@400&display=swap" rel="stylesheet"/>
<link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/katex@0.16.11/dist/katex.min.css"/>
<script defer src="https://cdn.jsdelivr.net/npm/katex@0.16.11/dist/katex.min.js"></script>
<script defer src="https://cdn.jsdelivr.net/npm/katex@0.16.11/dist/contrib/auto-render.min.js"
  onload="renderMathInElement(document.body,{{delimiters:[{{left:'$$',right:'$$',display:true}},{{left:'$',right:'$',display:false}},{{left:'\\\\(',right:'\\\\)',display:false}},{{left:'\\\\[',right:'\\\\]',display:true}}]}});"></script>
<style>
*{{box-sizing:border-box;margin:0;padding:0}}
body{{background:#0a0a0a;color:#c0c0c0;font-family:'DM Sans',sans-serif;font-size:14px;line-height:1.7;padding:40px 20px}}
.wrap{{max-width:720px;margin:0 auto}}
.head{{margin-bottom:28px;padding-bottom:18px;border-bottom:1px solid rgba(255,255,255,.08)}}
h1{{font-size:20px;font-weight:600;color:#f0f0f0;margin-bottom:4px}}
.meta{{font-size:12px;color:#505050}}
.msg{{margin-bottom:20px;padding:14px 16px;border-radius:12px;border:1px solid rgba(255,255,255,.07)}}
.msg.assistant{{background:#111}}
.msg.user{{background:#161616;border-color:rgba(255,255,255,.05)}}
.label{{font-size:11px;font-weight:600;letter-spacing:.07em;text-transform:uppercase;color:#505050;margin-bottom:7px}}
.msg.assistant .label{{color:hsl(38,88%,55%)}}
.text pre{{background:#0d0d0d;border:1px solid rgba(255,255,255,.07);border-radius:7px;padding:12px;overflow-x:auto;margin:8px 0;font-family:'DM Mono',monospace;font-size:12px}}
.text code{{font-family:'DM Mono',monospace;font-size:12px}}
.badge{{display:inline-block;margin-top:24px;padding:4px 12px;border-radius:99px;font-size:11px;background:rgba(255,255,255,.05);color:#505050}}
</style></head>
<body><div class="wrap">
<div class="head"><h1>{title}</h1><p class="meta">Shared from VedAI · vedai.soumyamazumdar.com</p></div>
{msgs_html}
<div class="badge">Shared from VedAI</div>
</div></body></html>''')

# ── Entry ─────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    import uvicorn
    uvicorn.run("app:app", host=os.getenv("HOST", "127.0.0.1"),
                port=int(os.getenv("PORT", "8000")), reload=False)
